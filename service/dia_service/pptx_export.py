"""Export a diastil dialect deck (§profile) to a .pptx — the OUTBOUND bridge.

`ingest/pptx.ts` brings PowerPoint *into* the dialect; this module takes a
saved dialect deck the other way: dialect HTML -> a native `.pptx`. Uploaded
with the Google Slides mime type it becomes a real, editable Google Slides
deck (`gdrive upload --mime-type application/vnd.google-apps.presentation`),
so a diastil deck reaches anyone who lives in Slides.

Why a .pptx and not "print to Slides": a .pptx converts to NATIVE Slides
objects — every text box, shape, and connector stays editable in Slides,
re-themeable, shareable. Charts and scene diagrams are drawn as vector shapes
(rectangles, ovals, connectors, freeforms) rather than a rasterized chart
object, because Slides rasterizes imported chart *objects* (they blur) while
plain shapes stay crisp and editable.

Fidelity model, mirroring ingest's: the dialect already carries a derived,
geometry-true rendering (scene nodes carry `data-x/y/w/h`; edges carry their
routed `<path d>`; charts carry a derived group) — so scenes and charts map
FAITHFULLY by their own coordinates (viewBox units scale linearly into the
figure box). Text/flow layout is mapped SEMANTICALLY from the dialect's layout
containers (cover / columns / stack) and text roles to slide geometry: this
yields clean, on-brand, editable Slides rather than a pixel tracing.

The dialect design space is 1280x720 px (16:9); a px maps to EMU at
9525 EMU/px (914400/96), so the deck fills a 12192000x6858000 EMU slide.

Stdlib-only HTML parsing (html.parser); python-pptx is the only third-party
dep. Entry points: `deck_to_pptx(html) -> bytes`, `export_file(src, dst)`.
"""

from __future__ import annotations

from dataclasses import dataclass, field
from html.parser import HTMLParser
import io
import re

import pptx
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Pt

# ---------------------------------------------------------------------------
# geometry: the dialect design space -> EMU
# ---------------------------------------------------------------------------

EMU_PER_PX = 9525  # 914400 EMU/inch / 96 px/inch — matches ingest/pptx.ts
DESIGN_W = 1280
DESIGN_H = 720
SLIDE_W = Emu(DESIGN_W * EMU_PER_PX)  # 12192000
SLIDE_H = Emu(DESIGN_H * EMU_PER_PX)  # 6858000


def _emu(px: float) -> int:
  return int(round(px * EMU_PER_PX))


# px -> pt for fonts (72pt/in over 96px/in).
def _pt(px: float) -> float:
  return px * 0.75


def _font_pt(size: float) -> "Pt":
  """python-pptx rejects font sizes outside 1pt..4000pt (an imported foreign

  deck can carry theme tokens that resolve to a sub-1pt size). Clamp so a weird
  input never 500s the export.
  """
  try:
    s = float(size)
  except (TypeError, ValueError):
    s = 12.0
  return Pt(max(1.0, min(4000.0, s)))


# ---------------------------------------------------------------------------
# a tiny DOM (stdlib) — enough to walk the dialect
# ---------------------------------------------------------------------------

VOID = {
    "area",
    "base",
    "br",
    "col",
    "embed",
    "hr",
    "img",
    "input",
    "link",
    "meta",
    "param",
    "source",
    "track",
    "wbr",
}


@dataclass
class El:
  tag: str
  attrs: dict[str, str]
  parent: "El | None" = None
  children: list["El"] = field(default_factory=list)
  text: str = ""  # direct text content (concatenated)

  def classes(self) -> set[str]:
    return set((self.attrs.get("class") or "").split())

  def has(self, cls: str) -> bool:
    return cls in self.classes()

  def find(self, pred) -> "El | None":
    for el in self.walk():
      if el is not self and pred(el):
        return el
    return None

  def find_all(self, pred) -> list["El"]:
    return [el for el in self.walk() if el is not self and pred(el)]

  def walk(self):
    yield self
    for c in self.children:
      yield from c.walk()

  def all_text(self) -> str:
    parts = []
    for el in self.walk():
      if el.text:
        parts.append(el.text)
    return re.sub(r"\s+", " ", " ".join(parts)).strip()


class _Tree(HTMLParser):

  def __init__(self) -> None:
    super().__init__(convert_charrefs=True)
    self.root = El("#root", {})
    self._stack = [self.root]

  def handle_starttag(self, tag, attrs):
    el = El(tag, {k: v or "" for k, v in attrs}, parent=self._stack[-1])
    self._stack[-1].children.append(el)
    if tag not in VOID:
      self._stack.append(el)

  def handle_startendtag(self, tag, attrs):
    el = El(tag, {k: v or "" for k, v in attrs}, parent=self._stack[-1])
    self._stack[-1].children.append(el)

  def handle_endtag(self, tag):
    for i in range(len(self._stack) - 1, 0, -1):
      if self._stack[i].tag == tag:
        del self._stack[i:]
        break

  def handle_data(self, data):
    if data.strip():
      self._stack[-1].text = (self._stack[-1].text + data).strip()


def _parse(html: str) -> El:
  t = _Tree()
  t.feed(html)
  return t.root


# ---------------------------------------------------------------------------
# theme
# ---------------------------------------------------------------------------

_DEFAULTS = {
    "paper": "#ffffff",
    "ink": "#1a1a1a",
    "ink-soft": "#33312b",
    "ink-faint": "#85837a",
    "accent": "#1e6fcc",
    "rule": "#c6c3b6",
}
_SCALE_DEFAULT = {1: 12.0, 2: 15.0, 3: 18.0, 4: 22.0, 5: 30.0, 6: 38.0, 7: 48.0}


@dataclass
class Theme:
  colors: dict[str, str]
  scale: dict[int, float]
  face_display: str = "Arial"
  face_body: str = "Arial"
  face_label: str = "Consolas"
  pad: float = 52.0
  gap: float = 24.0

  def rgb(self, key: str, fallback: str = "ink") -> RGBColor:
    return _hex(self.colors.get(key) or self.colors.get(fallback) or "#1a1a1a")


def _hex(s: str) -> RGBColor:
  s = (s or "").strip().lstrip("#")
  if len(s) == 3:
    s = "".join(c * 2 for c in s)
  if len(s) != 6 or not re.fullmatch(r"[0-9a-fA-F]{6}", s):
    return RGBColor.from_string("1a1a1a")
  return RGBColor.from_string(s.lower())


_FIRST_FONT = re.compile(r'"([^"]+)"|\'([^\']+)\'|([^,]+)')


def _first_font(css: str) -> str:
  m = _FIRST_FONT.search(css or "")
  if not m:
    return "Arial"
  return (m.group(1) or m.group(2) or m.group(3) or "Arial").strip()


def _read_theme(root: El) -> Theme:
  style = root.find(
      lambda e: e.tag == "style" and e.attrs.get("id") == "dia-theme"
  )
  css = style.all_text() if style else ""
  # Custom properties live in :root { --dia-*: v; }
  props: dict[str, str] = {}
  for m in re.finditer(r"--dia-([a-z0-9-]+)\s*:\s*([^;}\n]+)", css, re.I):
    props[m.group(1).strip().lower()] = m.group(2).strip()

  colors = dict(_DEFAULTS)
  for k in list(colors):
    if k in props:
      colors[k] = props[k]
  scale = dict(_SCALE_DEFAULT)
  for i in range(1, 8):
    v = props.get(f"scale-{i}")
    if v:
      mm = re.search(r"(-?\d+(?:\.\d+)?)", v)
      if mm:
        scale[i] = float(mm.group(1))

  def _px(name: str, default: float) -> float:
    v = props.get(name)
    if v:
      mm = re.search(r"(-?\d+(?:\.\d+)?)", v)
      if mm:
        return float(mm.group(1))
    return default

  return Theme(
      colors=colors,
      scale=scale,
      face_display=_first_font(props.get("face-display", "")),
      face_body=_first_font(props.get("face-body", "")),
      face_label=_first_font(props.get("face-label", "")),
      pad=_px("pad", 52.0),
      gap=_px("gap", 24.0),
  )


def _resolve_color(theme: Theme, value: str | None, fallback: str) -> RGBColor:
  """A per-node/edge style value: `var(--dia-rule)` or a hex literal."""
  if not value:
    return theme.rgb(fallback)
  m = re.search(r"var\(\s*--dia-([a-z0-9-]+)", value, re.I)
  if m:
    return theme.rgb(m.group(1).lower(), fallback)
  lit = re.search(r"#[0-9a-f]{3,6}", value, re.I)
  if lit:
    return _hex(lit.group(0))
  return theme.rgb(fallback)


def _style_prop(el: El, name: str) -> str | None:
  st = el.attrs.get("style") or ""
  m = re.search(re.escape(name) + r"\s*:\s*([^;]+)", st, re.I)
  return m.group(1).strip() if m else None


# ---------------------------------------------------------------------------
# low-level drawing helpers
# ---------------------------------------------------------------------------


def _fill(
    shape,
    color: RGBColor | None,
    line: RGBColor | None = None,
    line_w: float = 0.0,
) -> None:
  if color is None:
    shape.fill.background()
  else:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
  if line is None:
    shape.line.fill.background()
  else:
    shape.line.color.rgb = line
    shape.line.width = Pt(line_w or 1.0)
  shape.shadow.inherit = False


def _bg(slide, theme: Theme) -> None:
  r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
  _fill(r, theme.rgb("paper"))
  sp = r._element
  sp.getparent().remove(sp)
  slide.shapes._spTree.insert(2, sp)


def _textbox(
    slide,
    x,
    y,
    w,
    h,
    runs,
    *,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
    leading: float = 1.2,
    space_after=6.0,
):
  """runs: list[paragraph]; paragraph: list[(text, pt, color, bold, font)]."""
  tb = slide.shapes.add_textbox(_emu(x), _emu(y), _emu(w), _emu(h))
  tf = tb.text_frame
  tf.word_wrap = True
  tf.vertical_anchor = anchor
  for i, para in enumerate(runs):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    p.line_spacing = leading
    for txt, size, color, bold, font in para:
      r = p.add_run()
      r.text = txt
      r.font.size = _font_pt(size)
      r.font.color.rgb = color
      r.font.bold = bold
      r.font.name = font
  return tb


# ---------------------------------------------------------------------------
# scene (diagram) rendering — faithful, from the derived geometry
# ---------------------------------------------------------------------------

_SHAPE_MAP = {
    "rect": MSO_SHAPE.RECTANGLE,
    "rounded": MSO_SHAPE.ROUNDED_RECTANGLE,
    "pill": MSO_SHAPE.ROUNDED_RECTANGLE,
    "ellipse": MSO_SHAPE.OVAL,
    "diamond": MSO_SHAPE.DIAMOND,
    "cylinder": MSO_SHAPE.CAN,
    "hex": MSO_SHAPE.HEXAGON,
    "parallelogram": MSO_SHAPE.PARALLELOGRAM,
    "triangle": MSO_SHAPE.ISOSCELES_TRIANGLE,
    "cloud": MSO_SHAPE.CLOUD,
    "note": MSO_SHAPE.FOLDED_CORNER,
}


def _viewbox(svg: El) -> tuple[float, float, float, float]:
  vb = svg.attrs.get("viewbox") or svg.attrs.get("viewBox") or "0 0 340 250"
  nums = [float(x) for x in re.findall(r"-?\d+(?:\.\d+)?", vb)][:4]
  while len(nums) < 4:
    nums.append([0, 0, 340, 250][len(nums)])
  return nums[0], nums[1], nums[2], nums[3]


def _fit(vw: float, vh: float, box: tuple[float, float, float, float]):
  """Return a mapper (vx,vy)->(px,py) fitting viewBox into a px box, centered.

  Guards against a degenerate box (a caller whose text overflowed can hand us a
  non-positive width/height); a negative scale would flip geometry off-slide, so
  clamp to a zero-scale no-op mapper collapsed at the box origin instead.
  """
  bx, by, bw, bh = box
  if vw <= 0 or vh <= 0 or bw <= 0 or bh <= 0:
    return (lambda x, y: (bx, by)), 0.0
  scale = min(bw / vw, bh / vh)
  rw, rh = vw * scale, vh * scale
  ox, oy = bx + (bw - rw) / 2, by + (bh - rh) / 2

  def m(x: float, y: float) -> tuple[float, float]:
    return ox + x * scale, oy + y * scale

  return m, scale


# One SVG number (leading-dot `.5`, trailing-dot `10.`, exponent `1e2` all legal)
# OR a path command letter. `\d*\.\d+` precedes `\d+\.?` so `10.5` matches whole.
_PATH_TOKEN = re.compile(
    r"[MmLlHhVvCcSsQqTtAaZz]|[-+]?(?:\d*\.\d+|\d+\.?)(?:[eE][-+]?\d+)?"
)


def _path_points(d: str) -> list[tuple[float, float]]:
  """Approximate a path's `d` as a polyline (M/L/H/V + bezier endpoints).

  Good enough for connectors: dialect edge paths are straight, orthogonal
  (H/V), or a single quadratic/cubic we sample at its endpoints + midpoint.
  """
  toks = _PATH_TOKEN.findall(d or "")
  pts: list[tuple[float, float]] = []
  i, cx, cy, cmd = 0, 0.0, 0.0, ""

  def num() -> float:
    nonlocal i
    if i >= len(toks):  # truncated path (too few coords) — don't IndexError
      return 0.0
    v = float(toks[i])
    i += 1
    return v

  while i < len(toks):
    t = toks[i]
    if t.isalpha():
      cmd = t
      i += 1
      if cmd in "Zz":
        continue
    rel = cmd.islower()
    c = cmd.upper()
    if c == "M" or c == "L":
      x, y = num(), num()
      cx, cy = (cx + x, cy + y) if rel else (x, y)
      pts.append((cx, cy))
    elif c == "H":
      x = num()
      cx = cx + x if rel else x
      pts.append((cx, cy))
    elif c == "V":
      y = num()
      cy = cy + y if rel else y
      pts.append((cx, cy))
    elif c == "Q":
      x1, y1, x, y = num(), num(), num(), num()
      ex, ey = (cx + x, cy + y) if rel else (x, y)
      mx, my = (cx + x1, cy + y1) if rel else (x1, y1)
      pts.append((mx, my))
      pts.append((ex, ey))
      cx, cy = ex, ey
    elif c == "C":
      x1, y1, x2, y2, x, y = (num() for _ in range(6))
      ex, ey = (cx + x, cy + y) if rel else (x, y)
      m2x, m2y = (cx + x2, cy + y2) if rel else (x2, y2)
      pts.append((m2x, m2y))
      pts.append((ex, ey))
      cx, cy = ex, ey
    else:  # unsupported command token stream — bail with what we have
      i += 1
  return pts


def _draw_node(slide, theme: Theme, node: El, m, scale: float) -> None:
  def fnum(k: str, d: float) -> float:
    try:
      return float(node.attrs.get(k, d))
    except (TypeError, ValueError):
      return d

  x, y = fnum("data-x", 0), fnum("data-y", 0)
  w, h = fnum("data-w", 120), fnum("data-h", 40)
  shape_name = (node.attrs.get("data-shape") or "rounded").lower()
  emphasis = "data-dia-emphasis" in node.attrs

  px, py = m(x, y)
  pw, ph = w * scale, h * scale
  fill = _resolve_color(theme, _style_prop(node, "--dia-node-fill"), "paper")
  stroke = _resolve_color(
      theme,
      _style_prop(node, "--dia-node-stroke"),
      "accent" if emphasis else "ink",
  )
  try:
    sw = float(
        _style_prop(node, "--dia-node-stroke-w") or (2.0 if emphasis else 1.3)
    )
  except (TypeError, ValueError):
    sw = 2.0 if emphasis else 1.3

  sp = None
  data_path = node.attrs.get("data-path")
  if shape_name == "path" and data_path:
    # Freeform node: data-path is a 0..100-normalized outline scaled into the
    # node box. Compose normalized -> viewBox -> slide so it lands where the
    # editor draws it (a star, ring, blob stays that shape, not a rectangle).
    def mp(nx: float, ny: float) -> tuple[float, float]:
      return m(x + nx / 100.0 * w, y + ny / 100.0 * h)

    _freeform(
        slide,
        theme,
        _path_points(data_path),
        mp,
        fill,
        stroke,
        sw,
        scale,
        close=True,
    )
  else:
    mso = _SHAPE_MAP.get(shape_name, MSO_SHAPE.ROUNDED_RECTANGLE)
    sp = slide.shapes.add_shape(mso, _emu(px), _emu(py), _emu(pw), _emu(ph))
    _fill(sp, fill, stroke, _line_pt(sw, scale))

  label = node.attrs.get("data-label")
  if not label:
    lab = node.find(lambda e: e.has("dia-node-label"))
    label = lab.all_text() if lab else ""
  if label:
    ink = _resolve_color(theme, _style_prop(node, "--dia-node-ink"), "ink")
    if sp is not None:
      tf = sp.text_frame
      tf.word_wrap = True
      tf.vertical_anchor = MSO_ANCHOR.MIDDLE
      p = tf.paragraphs[0]
      p.alignment = PP_ALIGN.CENTER
      r = p.add_run()
      r.text = label
      r.font.size = _font_pt(_pt(12 * scale + 4))
      r.font.color.rgb = ink
      r.font.name = theme.face_body
    else:  # freeform node — center a label textbox over the node box
      _textbox(
          slide,
          px,
          py,
          pw,
          ph,
          [[(label, _pt(12 * scale + 4), ink, False, theme.face_body)]],
          align=PP_ALIGN.CENTER,
          anchor=MSO_ANCHOR.MIDDLE,
      )


def _draw_edge(slide, theme: Theme, edge: El, m, scale: float) -> None:
  path = edge.find(lambda e: e.tag == "path" and e.has("dia-edge-path"))
  stroke = _resolve_color(theme, _style_prop(edge, "--dia-edge-stroke"), "ink")
  try:
    ew = float(_style_prop(edge, "--dia-edge-w") or 1.2)
  except (TypeError, ValueError):
    ew = 1.2

  pts = _path_points(path.attrs.get("d", "")) if path else []
  if len(pts) >= 2:
    mapped = [m(x, y) for x, y in pts]
    # Freeform polyline through the routed points (keeps ortho bends/curves).
    fb = slide.shapes.build_freeform(
        _emu(mapped[0][0]), _emu(mapped[0][1]), scale=1.0
    )
    fb.add_line_segments(
        [(_emu(x), _emu(y)) for x, y in mapped[1:]], close=False
    )
    shp = fb.convert_to_shape()
    shp.fill.background()
    shp.line.color.rgb = stroke
    shp.line.width = Pt(ew)
    shp.shadow.inherit = False
    _arrow_end(shp)

  # Label: prefer the derived <text class="dia-edge-label" x y> — it carries the
  # router's own placement, so shared-corridor labels don't collide. Fall back
  # to data-label at the path midpoint only when there is no derived text.
  lab = edge.find(lambda e: e.has("dia-edge-label"))
  label = (lab.all_text() if lab else "") or edge.attrs.get("data-label") or ""
  if label:
    lx: float | None = None
    ly: float | None = None
    xs = lab.attrs.get("x") if lab is not None else None
    ys = lab.attrs.get("y") if lab is not None else None
    if xs is not None and ys is not None:
      try:
        lx, ly = m(float(xs), float(ys))
      except ValueError:
        lx = ly = None
    if lx is None and pts:
      mid = pts[len(pts) // 2]
      lx, ly = m(mid[0], mid[1])
    if lx is not None and ly is not None:
      _textbox(
          slide,
          lx - 60,
          ly - 16,
          120,
          20,
          [[(
              label,
              _pt(11 * scale + 3),
              theme.rgb("ink-soft"),
              False,
              theme.face_label,
          )]],
          align=PP_ALIGN.CENTER,
      )


def _arrow_end(shape) -> None:
  """Set a triangular arrowhead on the line end via the drawingml XML."""
  from pptx.oxml.ns import qn

  ln = shape.line._get_or_add_ln()
  tail = ln.find(qn("a:tailEnd"))
  if tail is None:
    tail = ln.makeelement(qn("a:tailEnd"), {})
    ln.append(tail)
  tail.set("type", "triangle")
  tail.set("w", "med")
  tail.set("len", "med")


# ---------------------------------------------------------------------------
# generic inline-SVG primitives — decorative art, brand marks, freeform strokes
# ---------------------------------------------------------------------------

_DRAW_TAGS = {
    "path",
    "circle",
    "ellipse",
    "rect",
    "line",
    "polyline",
    "polygon",
    "text",
}


def _svg_paint(theme: Theme, el: El, attr: str, default: str | None):
  """Resolve an svg fill/stroke: token var(), hex, currentColor, or none."""
  v = _style_prop(el, attr) or el.attrs.get(attr)
  if v is None:
    return theme.rgb(default) if default else None
  v = v.strip().lower()
  if v in ("none", "transparent"):
    return None
  tok = re.search(r"--dia-([a-z0-9-]+)", v)
  if tok:
    return theme.rgb(tok.group(1))
  lit = re.search(r"#[0-9a-f]{3,8}", v)
  if lit:
    return _hex(lit.group(0)[:7])
  if v == "currentcolor":
    return theme.rgb("ink")
  return theme.rgb(default) if default else None


def _svg_f(el: El, name: str, d: float = 0.0) -> float:
  try:
    return float(el.attrs.get(name))  # type: ignore[arg-type]
  except (TypeError, ValueError):
    return d


def _svg_stroke_w(el: El, d: float = 1.0) -> float:
  v = _style_prop(el, "stroke-width") or el.attrs.get("stroke-width")
  try:
    return float(v)  # type: ignore[arg-type]
  except (TypeError, ValueError):
    return d


def _svg_font_px(el: El, d: float) -> float:
  st = el.attrs.get("style", "") or ""
  mm = re.search(r"font(?:-size)?:\s*(\d+(?:\.\d+)?)px", st)
  return float(mm.group(1)) if mm else d


def _line_pt(sw: float, scale: float) -> float:
  return max(0.5, _pt(sw * scale))


def _freeform(slide, theme, coords, m, fill, stroke, sw, scale, close) -> None:
  if len(coords) < 2:
    return
  mp = [m(x, y) for x, y in coords]
  fb = slide.shapes.build_freeform(_emu(mp[0][0]), _emu(mp[0][1]), scale=1.0)
  fb.add_line_segments([(_emu(x), _emu(y)) for x, y in mp[1:]], close=close)
  shp = fb.convert_to_shape()
  if fill is not None and close:
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
  else:
    shp.fill.background()
  if stroke is not None:
    shp.line.color.rgb = stroke
    shp.line.width = Pt(_line_pt(sw, scale))
  elif fill is None:
    shp.line.color.rgb = theme.rgb("ink")
    shp.line.width = Pt(_line_pt(sw, scale))
  else:
    shp.line.fill.background()
  shp.shadow.inherit = False


def _draw_prim(slide, theme: Theme, el: El, m, scale: float) -> None:
  """Render one svg drawable primitive as a native shape/connector/textbox."""
  tag = el.tag
  fill = _svg_paint(theme, el, "fill", None)
  stroke = _svg_paint(theme, el, "stroke", None)
  sw = _svg_stroke_w(el, 1.2)
  if tag in ("circle", "ellipse"):
    cx, cy = _svg_f(el, "cx"), _svg_f(el, "cy")
    rx = _svg_f(el, "r") or _svg_f(el, "rx")
    ry = _svg_f(el, "r") or _svg_f(el, "ry")
    if fill is None and stroke is None:
      fill = theme.rgb("ink")
    x, y = m(cx - rx, cy - ry)
    sp = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        _emu(x),
        _emu(y),
        _emu(2 * rx * scale),
        _emu(2 * ry * scale),
    )
    _fill(sp, fill, stroke, _line_pt(sw, scale))
  elif tag == "rect":
    x0, y0 = _svg_f(el, "x"), _svg_f(el, "y")
    if fill is None and stroke is None:
      fill = theme.rgb("ink")
    x, y = m(x0, y0)
    sp = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        _emu(x),
        _emu(y),
        _emu(_svg_f(el, "width") * scale),
        _emu(_svg_f(el, "height") * scale),
    )
    _fill(sp, fill, stroke, _line_pt(sw, scale))
  elif tag == "line":
    p1 = m(_svg_f(el, "x1"), _svg_f(el, "y1"))
    p2 = m(_svg_f(el, "x2"), _svg_f(el, "y2"))
    cn = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        _emu(p1[0]),
        _emu(p1[1]),
        _emu(p2[0]),
        _emu(p2[1]),
    )
    cn.line.color.rgb = stroke or theme.rgb("ink")
    cn.line.width = Pt(_line_pt(sw, scale))
    cn.shadow.inherit = False
  elif tag in ("polyline", "polygon"):
    nums = [
        float(n)
        for n in re.findall(
            r"-?\d+(?:\.\d+)?", el.attrs.get("points", "") or ""
        )
    ]
    coords = [(nums[i], nums[i + 1]) for i in range(0, len(nums) - 1, 2)]
    _freeform(
        slide,
        theme,
        coords,
        m,
        fill,
        stroke,
        sw,
        scale,
        close=(tag == "polygon"),
    )
  elif tag == "path":
    d = el.attrs.get("d", "") or ""
    coords = _path_points(d)
    _freeform(
        slide,
        theme,
        coords,
        m,
        fill,
        stroke,
        sw,
        scale,
        close=("z" in d.lower()) or fill is not None,
    )
  elif tag == "text":
    txt = el.all_text()
    if not txt:
      return
    fs = _svg_font_px(el, 12.0)
    x, y = m(_svg_f(el, "x"), _svg_f(el, "y"))
    _textbox(
        slide,
        x - 4,
        y - fs * scale,
        420,
        fs * scale * 1.6 + 8,
        [[(
            txt,
            _pt(fs * scale),
            fill or theme.rgb("ink-soft"),
            False,
            theme.face_body,
        )]],
    )


def _ancestor_matches(el: El, pred) -> bool:
  p = el.parent
  while p is not None:
    if pred(p):
      return True
    p = p.parent
  return False


def _draw_prims(
    slide, theme: Theme, container: El, m, scale: float, skip=None
) -> None:
  for el in container.walk():
    if el is container or el.tag not in _DRAW_TAGS:
      continue
    if skip and skip(el):
      continue
    try:
      _draw_prim(slide, theme, el, m, scale)
    except (ValueError, TypeError, IndexError, ZeroDivisionError, KeyError):
      pass  # one malformed primitive must never break the deck


def _svg_declared_size(svg: El) -> tuple[float, float]:
  """On-slide px size of a decorative svg: its style width/height, else a

  viewBox-derived default.
  """
  st = svg.attrs.get("style", "") or ""

  def px(name: str) -> float:
    mm = re.search(name + r":\s*(\d+(?:\.\d+)?)px", st)
    return float(mm.group(1)) if mm else 0.0

  _, _, vw, vh = _viewbox(svg)
  w, h = px("width"), px("height")
  ratio = (vh / vw) if vw else 0.2
  if w and not h:
    h = w * ratio
  elif h and not w:
    w = h / ratio if ratio else h * 5
  elif not w and not h:
    w, h = 120.0, 120.0 * ratio
  return w, h


def _draw_decorative(slide, theme: Theme, svg: El, x: float, y: float) -> float:
  """Render a standalone decorative svg (brand mark, rule art) at (x,y) at its

  declared size; return its height so the caller can advance the cursor.
  """
  w, h = _svg_declared_size(svg)
  _, _, vw, vh = _viewbox(svg)
  m, scale = _fit(vw, vh, (x, y, w, h))
  _draw_prims(
      slide,
      theme,
      svg,
      m,
      scale,
      skip=lambda e: _ancestor_matches(
          e, lambda p: p.tag in ("defs", "marker")
      ),
  )
  return h


def _draw_scene(slide, theme: Theme, svg: El, box) -> None:
  _, _, vw, vh = _viewbox(svg)
  m, scale = _fit(vw, vh, box)
  # edges first (under nodes), then nodes. Each is isolated: one malformed
  # derived path/geometry skips that shape, never aborts the whole deck.
  for edge in svg.find_all(lambda e: "data-dia-edge" in e.attrs):
    try:
      _draw_edge(slide, theme, edge, m, scale)
    except (ValueError, TypeError, IndexError, ZeroDivisionError, KeyError):
      pass
  for node in svg.find_all(lambda e: "data-dia-node" in e.attrs):
    try:
      _draw_node(slide, theme, node, m, scale)
    except (ValueError, TypeError, IndexError, ZeroDivisionError, KeyError):
      pass
  # loose decorative primitives (freeform strokes, standalone labels) that live
  # directly in the scene — not inside a node/edge group or <defs>/<marker>.
  _draw_prims(
      slide,
      theme,
      svg,
      m,
      scale,
      skip=lambda e: _ancestor_matches(
          e,
          lambda p: "data-dia-node" in p.attrs
          or "data-dia-edge" in p.attrs
          or p.tag in ("defs", "marker"),
      ),
  )


# ---------------------------------------------------------------------------
# chart rendering — vector, from data attributes
# ---------------------------------------------------------------------------


def _parse_values(raw: str) -> list[tuple[str, float]]:
  out: list[tuple[str, float]] = []
  for part in re.split(r"[,;]", raw or ""):
    if ":" not in part:
      continue
    label, _, num = part.partition(":")
    try:
      out.append((label.strip(), float(num.strip())))
    except ValueError:
      pass
  return out


def _draw_chart(slide, theme: Theme, svg: El, box) -> None:
  kind = (svg.attrs.get("data-chart") or "bar").lower()
  data = _parse_values(svg.attrs.get("data-values") or "")
  if not data:
    return
  bx, by, bw, bh = box
  unit = svg.attrs.get("data-unit") or ""
  try:
    vmax = float(svg.attrs.get("data-max") or 0) or max(v for _, v in data)
  except (TypeError, ValueError):
    vmax = max(v for _, v in data) or 1.0
  vmax = vmax or 1.0

  plot_t = by + 10
  plot_h = bh - 48
  base_y = plot_t + plot_h
  accent = theme.rgb("accent")

  base = slide.shapes.add_shape(
      MSO_SHAPE.RECTANGLE, _emu(bx), _emu(base_y), _emu(bw), _emu(1.2)
  )
  _fill(base, theme.rgb("rule"))

  n = len(data)
  if kind in ("bar", "column"):
    slot = bw / n
    bar_w = slot * 0.6
    for i, (label, v) in enumerate(data):
      bh_px = (v / vmax) * plot_h if v > 0 else 0
      cx = bx + i * slot + (slot - bar_w) / 2
      if bh_px > 0:
        rect = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            _emu(cx),
            _emu(base_y - bh_px),
            _emu(bar_w),
            _emu(bh_px),
        )
        _fill(rect, accent)
      _textbox(
          slide,
          cx - 20,
          base_y - bh_px - 22,
          bar_w + 40,
          18,
          [[(
              _fmt(v) + unit,
              _pt(13),
              theme.rgb("ink"),
              True,
              theme.face_label,
          )]],
          align=PP_ALIGN.CENTER,
      )
      _textbox(
          slide,
          bx + i * slot,
          base_y + 6,
          slot,
          22,
          [[(label, _pt(13), theme.rgb("ink-soft"), False, theme.face_body)]],
          align=PP_ALIGN.CENTER,
      )
  else:  # line / scatter -> points + connecting polyline
    slot = bw / max(1, n - 1) if n > 1 else bw
    pts = []
    for i, (_, v) in enumerate(data):
      cx = bx + (i * slot if n > 1 else bw / 2)
      cy = base_y - (v / vmax) * plot_h
      pts.append((cx, cy))
    if kind == "line" and len(pts) >= 2:
      fb = slide.shapes.build_freeform(_emu(pts[0][0]), _emu(pts[0][1]))
      fb.add_line_segments(
          [(_emu(x), _emu(y)) for x, y in pts[1:]], close=False
      )
      shp = fb.convert_to_shape()
      shp.fill.background()
      shp.line.color.rgb = accent
      shp.line.width = Pt(2.0)
      shp.shadow.inherit = False
    for (cx, cy), (label, v) in zip(pts, data):
      dot = slide.shapes.add_shape(
          MSO_SHAPE.OVAL, _emu(cx - 4), _emu(cy - 4), _emu(8), _emu(8)
      )
      _fill(dot, accent)
      _textbox(
          slide,
          cx - 30,
          base_y + 6,
          60,
          22,
          [[(label, _pt(12), theme.rgb("ink-soft"), False, theme.face_body)]],
          align=PP_ALIGN.CENTER,
      )


def _fmt(v: float) -> str:
  return str(int(v)) if v == int(v) else f"{v:.1f}"


# ---------------------------------------------------------------------------
# slide layout — semantic mapping of dialect layout + roles to geometry
# ---------------------------------------------------------------------------


def _role_text(el: El) -> str:
  return el.all_text()


def _paras_of_body(body: El, theme: Theme):
  """A body region -> paragraphs (list of runs). <p>/<li> become lines."""
  blocks = body.find_all(lambda e: e.tag in ("p", "li"))
  runs = []
  for b in blocks:
    txt = b.all_text()
    if not txt:
      continue
    if b.tag == "li":
      runs.append([
          (
              "•  ",
              _pt(theme.scale[2]),
              theme.rgb("accent"),
              True,
              theme.face_body,
          ),
          (
              txt,
              _pt(theme.scale[2]),
              theme.rgb("ink-soft"),
              False,
              theme.face_body,
          ),
      ])
    else:
      runs.append([(
          txt,
          _pt(theme.scale[2]),
          theme.rgb("ink-soft"),
          False,
          theme.face_body,
      )])
  if not runs:
    txt = body.all_text()
    if txt:
      runs = [[(
          txt,
          _pt(theme.scale[2]),
          theme.rgb("ink-soft"),
          False,
          theme.face_body,
      )]]
  return runs


def _figure_of(slide_el: El) -> El | None:
  """The primary visual: a scene, chart, or image (in a figure or anywhere)."""
  return slide_el.find(
      lambda e: (e.tag == "svg" and (e.has("dia-scene") or e.has("dia-chart")))
      or e.tag == "img"
  )


def _draw_visual(slide, theme: Theme, vis: El, box) -> None:
  if vis.tag == "svg" and vis.has("dia-chart"):
    _draw_chart(slide, theme, vis, box)
  elif vis.tag == "svg":
    _draw_scene(slide, theme, vis, box)
  elif vis.tag == "img":
    _draw_image(slide, theme, vis, box)


def _draw_image(slide, theme: Theme, img: El, box) -> None:
  src = img.attrs.get("src") or ""
  bx, by, bw, bh = box
  m = re.match(r"data:image/(png|jpe?g|gif);base64,(.+)$", src, re.I | re.S)
  if m:
    import base64

    try:
      raw = base64.b64decode(m.group(2))
      slide.shapes.add_picture(
          io.BytesIO(raw), _emu(bx), _emu(by), _emu(bw), _emu(bh)
      )
      return
    except (ValueError, OSError):
      pass
  # SVG data-URI: decode + render its primitives as native vector shapes.
  svg_m = re.match(r"data:image/svg\+xml(;base64)?,(.+)$", src, re.I | re.S)
  if svg_m:
    import base64
    import urllib.parse

    payload = svg_m.group(2)
    try:
      markup = (
          base64.b64decode(payload).decode("utf-8")
          if svg_m.group(1)
          else urllib.parse.unquote(payload)
      )
      svg = _parse(markup).find(lambda e: e.tag == "svg")
      if svg is not None:
        _, _, vw, vh = _viewbox(svg)
        mm, scale = _fit(vw, vh, box)
        _draw_prims(
            slide,
            theme,
            svg,
            mm,
            scale,
            skip=lambda e: _ancestor_matches(
                e, lambda p: p.tag in ("defs", "marker")
            ),
        )
        return
    except (ValueError, OSError, UnicodeDecodeError):
      pass
  # remote refs can't embed directly — draw a labeled frame.
  frame = slide.shapes.add_shape(
      MSO_SHAPE.RECTANGLE, _emu(bx), _emu(by), _emu(bw), _emu(bh)
  )
  _fill(frame, theme.rgb("rule"), theme.rgb("ink-faint"), 1.0)
  _textbox(
      slide,
      bx,
      by + bh / 2 - 14,
      bw,
      28,
      [[(
          img.attrs.get("alt") or "figure",
          _pt(theme.scale[1]),
          theme.rgb("ink-soft"),
          False,
          theme.face_label,
      )]],
      align=PP_ALIGN.CENTER,
  )


def _add_notes(slide, slide_el: El) -> None:
  aside = slide_el.find(lambda e: e.tag == "aside" and e.has("dia-notes"))
  if aside:
    txt = aside.all_text()
    if txt:
      slide.notes_slide.notes_text_frame.text = txt


def _render_slide(prs, theme: Theme, slide_el: El) -> None:
  blank = prs.slide_layouts[6]
  slide = prs.slides.add_slide(blank)
  _bg(slide, theme)

  pad = theme.pad
  content_w = DESIGN_W - 2 * pad
  content_h = DESIGN_H - 2 * pad
  is_cover = slide_el.has("dia-cover")
  is_columns = slide_el.has("dia-columns")

  kicker = slide_el.find(lambda e: e.has("dia-kicker"))
  title = slide_el.find(lambda e: e.has("dia-title"))
  body = slide_el.find(lambda e: e.has("dia-body"))
  caption = slide_el.find(lambda e: e.has("dia-caption"))
  footnote = slide_el.find(lambda e: e.has("dia-footnote"))
  vis = _figure_of(slide_el)
  # A standalone decorative svg (brand mark, rule art) that is NOT the figure
  # scene/chart — rendered inline at the top of the head so it isn't dropped.
  deco = slide_el.find(
      lambda e: e.tag == "svg"
      and not e.has("dia-scene")
      and not e.has("dia-chart")
      and e is not vis
  )

  title_px = (
      theme.scale[7]
      if (title and title.has("dia-cover-title"))
      else theme.scale[5]
  )

  if is_columns and vis is not None:
    # 1.05fr / 1fr with a gap (mirrors .dia-columns).
    gap = theme.gap
    left_w = (content_w - gap) * (1.05 / 2.05)
    right_w = (content_w - gap) - left_w
    lx, rx = pad, pad + left_w + gap
    # left: text column, flowed
    y = pad
    y = _flow_head(slide, theme, kicker, title, title_px, lx, y, left_w)
    if body is not None:
      _textbox(
          slide,
          lx,
          y + 6,
          left_w,
          pad + content_h - y,
          _paras_of_body(body, theme),
          leading=1.5,
          space_after=10,
          anchor=MSO_ANCHOR.TOP,
      )
    # right: the visual, filling most of the column height
    vis_box = (rx, pad + 6, right_w, content_h - 40)
    _draw_visual(slide, theme, vis, vis_box)
    if caption is not None:
      _textbox(
          slide,
          rx,
          pad + content_h - 22,
          right_w,
          22,
          [[(
              caption.all_text(),
              _pt(theme.scale[1]),
              theme.rgb("ink-soft"),
              False,
              theme.face_label,
          )]],
      )
  elif is_cover:
    # vertically centered stack in the padded box
    stack = []
    if kicker is not None:
      stack.append(("kicker", kicker.all_text()))
    if title is not None:
      stack.append(("title", title.all_text()))
    if body is not None:
      stack.append(("body", body.all_text()))
    if caption is not None:
      stack.append(("caption", caption.all_text()))
    # estimate block heights to center
    est = 0.0
    for kind, txt in stack:
      est += _est_height(theme, kind, txt, content_w, title_px)
    deco_h = 0.0
    if deco is not None:
      _, deco_h = _svg_declared_size(deco)
      est += deco_h + 14  # brand mark + its margin-bottom
    y = pad + max(0, (content_h - est) / 2)
    if deco is not None:
      _draw_decorative(slide, theme, deco, pad, y)
      y += deco_h + 14
    y = _flow_head(slide, theme, kicker, title, title_px, pad, y, content_w)
    if body is not None:
      h = _est_height(theme, "body", body.all_text(), content_w, title_px)
      _textbox(
          slide,
          pad,
          y + 6,
          content_w,
          h + 20,
          _paras_of_body(body, theme),
          leading=1.5,
      )
      y += h + 12
    if caption is not None:
      _textbox(
          slide,
          pad,
          y + 10,
          content_w,
          30,
          [[(
              caption.all_text(),
              _pt(theme.scale[1]),
              theme.rgb("ink-soft"),
              False,
              theme.face_label,
          )]],
      )
  else:
    # default flow: head, then body, then any visual below
    y = pad
    y = _flow_head(slide, theme, kicker, title, title_px, pad, y, content_w)
    if body is not None:
      h = _est_height(theme, "body", body.all_text(), content_w, title_px)
      _textbox(
          slide,
          pad,
          y + 6,
          content_w,
          max(h + 20, 120),
          _paras_of_body(body, theme),
          leading=1.5,
          space_after=10,
      )
      y += h + 24
    if vis is not None:
      # Clamp: an overflowing head/body can push y past the content box; keep a
      # positive height so the visual renders (small) rather than collapsing.
      vis_h = max(80.0, pad + content_h - y - 10)
      _draw_visual(slide, theme, vis, (pad, y, content_w, vis_h))
    if caption is not None and vis is None:
      _textbox(
          slide,
          pad,
          pad + content_h - 22,
          content_w,
          22,
          [[(
              caption.all_text(),
              _pt(theme.scale[1]),
              theme.rgb("ink-soft"),
              False,
              theme.face_label,
          )]],
      )
  if footnote is not None:
    _textbox(
        slide,
        pad,
        DESIGN_H - pad + 4,
        content_w,
        20,
        [[(
            footnote.all_text(),
            _pt(theme.scale[1]),
            theme.rgb("ink-faint"),
            False,
            theme.face_label,
        )]],
    )
  _add_notes(slide, slide_el)


def _flow_head(slide, theme, kicker, title, title_px, x, y, w) -> float:
  if kicker is not None and kicker.all_text():
    _textbox(
        slide,
        x,
        y,
        w,
        24,
        [[(
            kicker.all_text().upper(),
            _pt(theme.scale[1]),
            theme.rgb("accent"),
            True,
            theme.face_label,
        )]],
    )
    y += theme.scale[1] + 12
  if title is not None and title.all_text():
    h = _est_height(theme, "title", title.all_text(), w, title_px)
    _textbox(
        slide,
        x,
        y,
        w,
        h + 10,
        [[(
            title.all_text(),
            _pt(title_px),
            theme.rgb("ink"),
            True,
            theme.face_display,
        )]],
        leading=1.14,
    )
    y += h + 14
  return y


def _est_height(
    theme: Theme, kind: str, text: str, w: float, title_px: float
) -> float:
  """Rough block height in px from text length + font size (for flow/centering)."""
  if kind == "title":
    size, lead, cpl = title_px, 1.14, max(8, w / (title_px * 0.56))
  elif kind == "kicker":
    return theme.scale[1] + 12
  elif kind == "caption":
    size, lead, cpl = theme.scale[1], 1.4, w / (theme.scale[1] * 0.55)
  else:  # body
    size, lead, cpl = theme.scale[2], 1.55, w / (theme.scale[2] * 0.52)
  lines = max(1, int(len(text) / max(1, cpl)) + 1)
  return lines * size * lead + 6


# ---------------------------------------------------------------------------
# public API
# ---------------------------------------------------------------------------


def deck_to_pptx(html: str) -> bytes:
  """Convert a dialect deck (HTML string) to .pptx bytes."""
  root = _parse(html)
  theme = _read_theme(root)
  slides = root.find_all(lambda e: e.tag == "section" and e.has("dia-slide"))
  prs = pptx.Presentation()
  prs.slide_width = SLIDE_W
  prs.slide_height = SLIDE_H
  for slide_el in slides:
    try:
      _render_slide(prs, theme, slide_el)
    except Exception:  # pylint: disable=broad-except
      # One malformed slide (common with imported foreign decks) must never
      # abort the whole export — emit a blank slide in its place and continue.
      try:
        prs.slides.add_slide(prs.slide_layouts[6])
      except Exception:  # pylint: disable=broad-except
        pass
  buf = io.BytesIO()
  prs.save(buf)
  return buf.getvalue()


def deck_title(html: str) -> str:
  root = _parse(html)
  t = root.find(lambda e: e.tag == "title")
  if t and t.all_text():
    return t.all_text()
  # first cover/slide title
  st = root.find(lambda e: e.has("dia-title"))
  return st.all_text() if st else "Presentation"


def deck_slide_count(html: str) -> int:
  """Number of `<section class="dia-slide">` in the deck (0 = not a deck)."""
  root = _parse(html)
  return len(root.find_all(lambda e: e.tag == "section" and e.has("dia-slide")))


def export_file(src: str, dst: str) -> int:
  with open(src, "r", encoding="utf-8") as f:
    html = f.read()
  data = deck_to_pptx(html)
  with open(dst, "wb") as f:
    f.write(data)
  root = _parse(html)
  n = len(root.find_all(lambda e: e.tag == "section" and e.has("dia-slide")))
  return n
